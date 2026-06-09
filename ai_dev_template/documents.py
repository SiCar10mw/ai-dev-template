"""Reusable document and presentation generation helpers."""

from __future__ import annotations

from collections.abc import Sequence
from pathlib import Path
from typing import Any
from zipfile import ZIP_DEFLATED, ZipFile, ZipInfo

from docx import Document
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

ACCENT = RGBColor(0x1F, 0x4E, 0x79)
INK = RGBColor(0x20, 0x20, 0x20)
LIGHT = RGBColor(0xF2, 0xF5, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

Section = tuple[str, str]


def build_docx(title: str, sections: Sequence[Section], output_path: Path) -> Path:
    """Build a simple DOCX document from section headings and body text."""
    output_path.parent.mkdir(parents=True, exist_ok=True)
    doc = Document()
    doc.add_heading(title, level=1)
    for heading, body in sections:
        doc.add_heading(heading, level=2)
        for paragraph in body.splitlines():
            if paragraph.strip():
                doc.add_paragraph(paragraph.strip())
    doc.save(str(output_path))
    _normalize_office_zip(output_path)
    return output_path


def build_reference_docx(output_path: Path) -> Path:
    """Build a generic reference DOCX that downstream projects can use as a brand template."""
    sections = [
        (
            "Purpose",
            "This file is a reusable reference document for generated project deliverables.",
        ),
        (
            "Documentation Rule",
            "Every code change must include an internal documentation update or a public docs-site update.",
        ),
    ]
    return build_docx("Project Reference Document", sections, output_path)


def build_deck(title: str, slides: Sequence[Section], output_path: Path) -> Path:
    """Build a concise PowerPoint deck from slide titles and body text."""
    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]

    title_slide = presentation.slides.add_slide(blank)
    _add_title_band(title_slide, title)
    _add_text(
        title_slide,
        "Generated from repository source.",
        Inches(0.75),
        Inches(2.15),
        Inches(11.5),
        Inches(0.8),
        24,
    )

    for slide_title, body in slides:
        slide = presentation.slides.add_slide(blank)
        _add_title_band(slide, slide_title)
        _add_body(slide, body)

    presentation.save(str(output_path))
    _normalize_office_zip(output_path)
    return output_path


def _normalize_office_zip(path: Path) -> None:
    """Normalize Office ZIP container timestamps so generated artifacts compare byte-for-byte."""
    fixed_date = (2026, 1, 1, 0, 0, 0)
    source_entries: list[tuple[str, bytes]] = []
    with ZipFile(path, "r") as source:
        for name in source.namelist():
            source_entries.append((name, source.read(name)))

    tmp_path = path.with_suffix(path.suffix + ".tmp")
    with ZipFile(tmp_path, "w", compression=ZIP_DEFLATED) as target:
        for name, data in source_entries:
            info = ZipInfo(name, fixed_date)
            info.compress_type = ZIP_DEFLATED
            target.writestr(info, data)
    tmp_path.replace(path)


def build_project_brief_deck(output_path: Path) -> Path:
    """Build a generic project brief deck."""
    slides = [
        ("Governance", "Constitution first\nHuman-gated external writes\nDeterministic authority for truth"),
        ("Delivery", "Spec-first delivery\nTest-driven implementation\nIndependent review before merge"),
        ("Documentation", "Docs impact is checked before commit\nPublic docs are curated in docs-site"),
    ]
    return build_deck("Project Brief", slides, output_path)


def _add_title_band(slide: Any, title: str) -> None:
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.05))
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.fill.background()
    _add_text(slide, title, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7), 30, WHITE, bold=True)


def _add_text(
    slide: Any,
    text: str,
    left: int,
    top: int,
    width: int,
    height: int,
    size: int,
    color: RGBColor = INK,
    *,
    bold: bool = False,
) -> None:
    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_body(slide: Any, body: str) -> None:
    textbox = slide.shapes.add_textbox(Inches(0.8), Inches(1.55), Inches(11.8), Inches(5.7))
    frame = textbox.text_frame
    frame.word_wrap = True
    for index, line in enumerate(body.splitlines()):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(8)
        run = paragraph.add_run()
        run.text = f"- {line.strip()}"
        run.font.size = Pt(21)
        run.font.color.rgb = INK
